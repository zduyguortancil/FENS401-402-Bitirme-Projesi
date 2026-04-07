"""
FENS401 Proje Sunumu PPTX Oluşturucu
"""
import subprocess, sys

# python-pptx kurulumu
subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Color palette ──
DARK_BG    = RGBColor(0x0F, 0x17, 0x2A)
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)
ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
ORANGE     = RGBColor(0xFF, 0x8C, 0x42)
GREEN      = RGBColor(0x06, 0xD6, 0xA0)

COVER_IMG = "/Users/zeynepduygu/.gemini/antigravity/brain/175dcf2d-92be-40a2-af5f-12c93b210a31/cover_airplane_1772043682281.png"

def add_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=18, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

# ═══════════════════════════════════════════════════
# SLIDE 1: KAPAK
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_dark_bg(slide)

# Cover image (full background)
if os.path.exists(COVER_IMG):
    slide.shapes.add_picture(COVER_IMG, Inches(0), Inches(0), W, H)

# Dark overlay
overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
# Make it semi-transparent via XML hack
from pptx.oxml.ns import qn
solidFill = overlay.fill._fill
alpha = solidFill.find(qn('a:solidFill'))
if alpha is not None:
    srgb = alpha.find(qn('a:srgbClr'))
    if srgb is not None:
        a_elem = srgb.makeelement(qn('a:alpha'), {'val': '55000'})
        srgb.append(a_elem)
overlay.line.fill.background()

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Flight Snapshot Dashboard", font_size=48, color=WHITE, bold=True)
add_accent_line(slide, Inches(1), Inches(3.1), Inches(3))
add_text_box(slide, Inches(1), Inches(3.4), Inches(10), Inches(0.8),
             "Havayolu Gelir Yönetimi & Talep Tahmin Sistemi", font_size=26, color=ACCENT2)
add_text_box(slide, Inches(1), Inches(5.5), Inches(10), Inches(0.6),
             "FENS 401 – Bitirme Projesi", font_size=20, color=LIGHT_GRAY)
add_text_box(slide, Inches(1), Inches(6.1), Inches(10), Inches(0.6),
             "2025 – 2026", font_size=18, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 2: PROJEMİZ NEDİR?
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             "Projemiz Nedir?", font_size=40, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.4), Inches(2.5))

add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
             "Havayolu şirketleri için uçuş bazında gelir ve doluluk verilerini analiz eden,\n"
             "gerçek zamanlı bir dashboard ve makine öğrenmesi tabanlı talep tahmin sistemi.",
             font_size=22, color=WHITE)

items = [
    "✈️  Uçuş bazında yolcu satış ve gelir takibi",
    "📊  Gerçek zamanlı dashboard ile anlık veri görselleştirme",
    "🤖  XGBoost tabanlı iki aşamalı talep tahmin modeli",
    "💰  Bilet geliri + ek gelir (ancillary) ayrıştırması",
    "🔍  Uçuş numarası, havalimanı ve tarih bazlı arama",
    "📈  Kabin sınıfı bazında doluluk oranı analizi",
]
add_bullet_slide(slide, Inches(0.8), Inches(3.2), Inches(11), Inches(4), items, font_size=22)

# ═══════════════════════════════════════════════════
# SLIDE 3: TEKNİK MİMARİ
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             "Teknik Mimari", font_size=40, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.4), Inches(2.5))

# Left column - Backend
add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5), Inches(0.6),
             "Backend", font_size=26, color=ORANGE, bold=True)
backend_items = [
    "🐍  Python + Flask web framework",
    "🦆  DuckDB – büyük veri sorguları (37M+ satır)",
    "📁  Parquet formatında veri depolama",
    "⚡  REST API endpointleri",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(3.5), backend_items, font_size=20)

# Right column - ML Model
add_text_box(slide, Inches(7), Inches(1.9), Inches(5), Inches(0.6),
             "Makine Öğrenmesi", font_size=26, color=GREEN, bold=True)
ml_items = [
    "🎯  İki aşamalı XGBoost modeli",
    "1️⃣   Classifier: Satış olacak mı?",
    "2️⃣   Regressor: Kaç yolcu satılacak?",
    "📋  31 özellik (feature) kullanımı",
]
add_bullet_slide(slide, Inches(7), Inches(2.6), Inches(5.5), Inches(3.5), ml_items, font_size=20)

# Bottom - Data
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.6),
             "Veri Seti", font_size=26, color=ACCENT2, bold=True)
data_items = [
    "📦  36,996,400 satır eğitim verisi  |  19 kolon  |  Economy & Business kabin sınıfları  |  DTD: 0–180 gün"
]
add_bullet_slide(slide, Inches(0.8), Inches(6.1), Inches(11), Inches(1.2), data_items, font_size=20)


# ═══════════════════════════════════════════════════
# SLIDE 4: NELER YAPTIK?
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             "Neler Yaptık?", font_size=40, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.4), Inches(2.5))

steps = [
    ("Sprint 1-2", "Veri Hazırlama & Altyapı", [
        "Uçuş snapshot ve metadata Parquet dosyaları oluşturuldu",
        "DuckDB ile bellek-dostu veri işleme pipeline'ı kuruldu",
        "V1 → V2 parquet geçişi (bilet + ek gelir ayrıştırması)",
    ]),
    ("Sprint 3", "Dashboard Geliştirme", [
        "Flask tabanlı web uygulaması geliştirildi",
        "Uçuş arama, tarih filtreleme ve kabin seçimi API'leri",
        "Gerçek zamanlı doluluk ve gelir grafikleri",
    ]),
    ("Sprint 4", "Talep Tahmin Modeli", [
        "37M satırlık eğitim veri seti hazırlandı",
        "XGBoost classifier + regressor eğitildi (31 feature)",
        "Forecast API entegrasyonu tamamlandı",
    ]),
]

y_pos = 1.8
for sprint_name, sprint_title, items in steps:
    add_text_box(slide, Inches(0.8), Inches(y_pos), Inches(2), Inches(0.5),
                 sprint_name, font_size=20, color=ORANGE, bold=True)
    add_text_box(slide, Inches(2.8), Inches(y_pos), Inches(4), Inches(0.5),
                 sprint_title, font_size=20, color=WHITE, bold=True)
    for j, item in enumerate(items):
        add_text_box(slide, Inches(2.8), Inches(y_pos + 0.45 + j * 0.4), Inches(9), Inches(0.4),
                     f"•  {item}", font_size=17, color=LIGHT_GRAY)
    y_pos += 0.45 + len(items) * 0.4 + 0.3

# ═══════════════════════════════════════════════════
# SLIDE 5: VERİ ANALİZİ SONUÇLARI
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             "Veri Analizi Sonuçları", font_size=40, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.4), Inches(2.5))

# Metric cards
metrics = [
    ("36.9M+", "Toplam Satır", ACCENT),
    ("19", "Özellik Sayısı", ORANGE),
    ("0.79", "Ort. Günlük Satış", GREEN),
    ("70.8%", "Sıfır Satış Oranı", RGBColor(0xFF, 0x61, 0x6B)),
]

card_width = Inches(2.7)
card_height = Inches(2.2)
gap = Inches(0.35)
start_x = Inches(0.8)

for i, (value, label, color) in enumerate(metrics):
    x = start_x + i * (card_width + gap)
    y = Inches(2)
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x3C)
    card.line.color.rgb = RGBColor(0x2A, 0x3A, 0x55)
    card.line.width = Pt(1)
    # Value
    add_text_box(slide, x, y + Inches(0.4), card_width, Inches(0.9),
                 value, font_size=42, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, x, y + Inches(1.4), card_width, Inches(0.5),
                 label, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Year distribution
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.6),
             "Yıl Dağılımı", font_size=22, color=WHITE, bold=True)
year_items = [
    "📅  2025: 18,484,806 satır   |   2026: 18,499,648 satır   |   2027: 11,946 satır",
    "✈️  Kabin Sınıfları: Economy (18.5M) + Business (18.5M) — dengeli dağılım",
    "📏  DTD Aralığı: 0 – 180 gün (kalkışa kalan gün sayısı)",
]
add_bullet_slide(slide, Inches(0.8), Inches(5.4), Inches(11), Inches(1.8), year_items, font_size=18)

# ═══════════════════════════════════════════════════
# SLIDE 6: DASHBOARD ÖZELLİKLERİ
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             "Dashboard Özellikleri", font_size=40, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.4), Inches(2.5))

features = [
    ("🔎  Akıllı Arama", "Uçuş numarası veya havalimanı koduna göre anlık arama"),
    ("📅  Tarih Filtresi", "Belirli bir tarihteki tüm uçuşları listeleme"),
    ("📊  Snapshot Görünümü", "Günlük yolcu satışı, kümülatif gelir, doluluk oranı"),
    ("💺  Kabin Analizi", "Economy ve Business ayrı ayrı analiz"),
    ("💰  Gelir Ayrıştırma", "Bilet geliri + ek gelir (ancillary) ayrı gösterim"),
    ("🤖  Talep Tahmini", "XGBoost modeli ile kalan günlerdeki beklenen talep"),
    ("📈  Trend Grafikleri", "DTD bazında satış eğilimi görselleştirme"),
]

for i, (title, desc) in enumerate(features):
    y = Inches(1.8) + i * Inches(0.72)
    add_text_box(slide, Inches(0.8), y, Inches(4), Inches(0.4),
                 title, font_size=21, color=ORANGE, bold=True)
    add_text_box(slide, Inches(4.5), y, Inches(8), Inches(0.4),
                 desc, font_size=19, color=WHITE)

# ═══════════════════════════════════════════════════
# SLIDE 7: TEŞEKKÜRLER
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

if os.path.exists(COVER_IMG):
    slide.shapes.add_picture(COVER_IMG, Inches(0), Inches(0), W, H)

overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, H)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
solidFill = overlay.fill._fill
alpha = solidFill.find(qn('a:solidFill'))
if alpha is not None:
    srgb = alpha.find(qn('a:srgbClr'))
    if srgb is not None:
        a_elem = srgb.makeelement(qn('a:alpha'), {'val': '65000'})
        srgb.append(a_elem)
overlay.line.fill.background()

add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
             "Teşekkürler", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5.5), Inches(4.1), Inches(2.3))
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
             "Sorularınız?", font_size=28, color=ACCENT2, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
             "FENS 401 – Bitirme Projesi  •  2025–2026", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = "/Users/zeynepduygu/Desktop/fens401/FENS401_Sunum.pptx"
prs.save(output_path)
print(f"\n✅ Sunum kaydedildi: {output_path}")
print(f"   Toplam slayt: {len(prs.slides)}")
